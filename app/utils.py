from pptx import Presentation
from openai import OpenAI
import os
from dotenv import load_dotenv
import tempfile
import shutil
from django.http import HttpResponse
from .models import PPTXDetails
from django.core.files.base import ContentFile




load_dotenv()  # take environment variables from .env.
api_key = os.getenv('API_KEY')
def generate_content(value):
    # Initialize OpenAI client
    client = OpenAI(api_key=api_key)

    # Generate content using GPT-3.5 model
    completion = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are a pptxassistant, skilled in creating note from large content with max 4 bulllet point"},
            {"role": "user", "content": f"{value}"}
        ]
    )

    return completion.choices[0].message.content

def generate_presentation(value , n = 2):
    # Create a presentation object
    prs = Presentation()

    # Slide 1: Title and Author Name
    slide_layout = prs.slide_layouts[0]  # Title Slide Layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = value
    for i in range(0, n):
        # Slide 2: Topic Heading and Content
        slide_layout = prs.slide_layouts[i]  # Content Slide Layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        content.text = generate_content(value)

    

    # Save the presentation to a temporary file
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
        prs.save(tmp_file.name)

        # Save the temporary file to the database
        with open(tmp_file.name, 'rb') as f:
            pptx_detail = PPTXDetails(pptx_file=ContentFile(f.read(), 'my_presentation.pptx'))
            pptx_detail.save()
            return pptx_detail
        # Clean up the temporary file
        
    
    

    

    # Clean up
    

   

